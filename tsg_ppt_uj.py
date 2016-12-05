#-*- coding: utf-8 -*-

###############TODO##############
#org - short + long version
#5 szabály tesztelés, kiiratni kieso orgokat level3on lesz ilyen
#test-4levelppts

#2. dia - táblázat helyett template
#egység fölötti összes szint eredménye + egy szintig az alatta levő eredmények (org-onként)
#százalékos arányok

#3.dia
#csak az adott organizáció eredménye minden kérdésre
#az előző évihez képest mennyit változott százalékban

#4.dia
#első kérdés minden szervezeti egységre (fölötte összes szint, alatta egy szint)
#előző évhez képest %-os változás a 4-5 válaszra
#1. az 1-2 válasz, 2. a 3. válasz, 3. a 4-5 válasz

#6-7 dia
#marad az 5 válasz
#egység fölötti összes szint, egység alatt egy szint

#####
##
#####
##
#####

###############/TODO#############
 

import codecs
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.util import Inches
from pptx.util import Pt
from pptx.util import Cm
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import time
import structure
import orgs
import csv
import sys
import os
import shutil

reload(sys)  
sys.setdefaultencoding('utf8')


#UTF8Writer = codecs.getwriter('utf8')



#box: default erteket allitunk be, ha negy parametert kap akkor nem lesz box, ha otot akkor hasznalja
#a=the number of the slide
#fill_slide_common(data, tsg_ppt, org_all, a, box=None)
#. az a vminek a vmije first_slidenak azt a shapejet keressuk ami a title (shape es title are fixed by pythonpptx)
def fill_slide_title(tsg_ppt, a, org_all):
  #datum=get_date
  first_slide=tsg_ppt.slides[0]
  common_slide=tsg_ppt.slides[a]
  org_1 = first_slide.placeholders[1]
  org = common_slide.placeholders[17]
  #text_frame = org.text_frame
  #text_frame.clear()
  #p = text_frame.paragraphs[0]
  #run = p.add_run()
  org.text=org_all
  org_1.text = org_all + "\n"+(time.strftime("%d.%m.%Y"))
  #subtitle = first_slide.placeholders[1] 
  #subtitle_2 = common_slide.placeholders[17]
  #subtitle.text=org_all
  return tsg_ppt
  #return -1: if the func doesnwork, we get with this parancs info about the wrong working, egyebkent exit and error code 

def fill_slide_common(tsg_ppt, org_all, a, c1, c2, c3, c4, c5, data, box=None):
#will be fill_slide_common(tsg_ppt, org_all, a):
  common_slide=tsg_ppt.slides[a]
  mittelwert = common_slide.placeholders[18]
  text_frame = mittelwert.text_frame
  text_frame.clear()
  p = text_frame.paragraphs[0]
  run = p.add_run()
  run.text=box
  font = run.font
  font.size=Pt(10.5)
  font.color.rgb = RGBColor(226, 0, 116)

  chart_data = ChartData()
  #chart_data.categories= [c1,c2,c3,c4,c5]
  chart_data.categories=['1a', '2a', 'a3', 'a4']
  #chart_data.add_series('01', (data[0], data[1], data[2], data[3], data[4]))
  a = [0.3,0.4,0.3] 
  b = [0.2,0.3, 0.5]
  c = [0.2,0.2,0.6]
  d = [0.1, 0.1, 0.8]
  #chart_data.add_series('1',(a,b,c))
  chart_data.add_series('2',(a[0], b[0], c[0], d[0]))
  chart_data.add_series('3',(a[1], b[1], c[1], d[1]))
  chart_data.add_series('4',(a[2], b[2], c[2], d[2]))
  x,y,cx,cy = Inches(0.3), Inches(2.25), Inches(9.38), Inches(4.7)
  graphic_frame = common_slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  chart.has_legend = True
  chart.plots[0].has_data_labels = True
  data_labels = chart.plots[0].data_labels
  data_labels.number_format = '0%'
  #nehasznaaalddata_labels.position = XL_LABEL_POSITION.OUTSIDE_END
  #value_axis.maximum_scale = 100.0
  tick_labels = value_axis.tick_labels
  tick_labels.number_format = '0%'
  tick_labels.font.size = Pt(12)
  tick_labels.font.type = 'Tele-GroteskEENor'
  chart.legend.position = XL_LEGEND_POSITION.BOTTOM
  chart.legend.font.size = Pt(12)
  chart.legend.font.type = 'Tele-GroteskEENor' 
  #plot = chart.plots[0]
  #plot.has_data_labels = True
  #data_labels = plot.data_labels
  data_labels.font.size = Pt(12)
  #data_labels.number_format = '0"%"'
  data_labels.font.color.rgb = RGBColor(0, 0, 0)
  #data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  #bar_plot = chart.plots[0]
  #bar_plot.gap_width = 20
  #bar_plot.overlap = -20
  
  chart.replace_data(chart_data)
  return tsg_ppt

def fill_table_slide(tsg_ppt,org_all, n):
  table_slide = tsg_ppt.slides[1]
  #rows = 3
  #cols = 2
  #left = Inches(3.62)
  #top = Inches(3.1)
  #width = Inches(5.0)
  #height = Inches (2.3)
  #table = table_slide.shapes.add_table(rows, cols, left, top, width, height).table
  # set column widths
  #table.columns[0].width = Inches(4.0)
  #table.columns[1].width = Inches(2.0)
  #table.cell(0, 0).text = 'Anzahl der Eingeladenen'
  #table.cell(0, 1).text = str(a)
  #table.cell(1, 0).text = 'Anzahl der Teilnehmer'
  #table.cell(1, 1).text = str(i)
  #table.cell(2, 0).text = 'Quote'
  #table.cell(2, 1).text = n
  chart_data = ChartData()
  chart_data.categories= [org_all]
  chart_data.add_series('01', (n))
  x,y,cx,cy = Inches(0.3), Inches(2.25), Inches(9.38), Inches(4.7)
  graphic_frame = table_slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  value_axis.maximum_scale = 100.0
  tick_labels = value_axis.tick_labels
  tick_labels.number_format = '0"%"'
  tick_labels.font.size = Pt(12)
  tick_labels.font.type = 'Tele-GroteskEENor'
  plot = chart.plots[0]
  plot.has_data_labels = True
  data_labels = plot.data_labels
  data_labels.font.size = Pt(12)
  data_labels.number_format = '0"%"'
  data_labels.font.color.rgb = RGBColor(0, 0, 0)
  data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  bar_plot = chart.plots[0]
  bar_plot.gap_width = 20
  bar_plot.overlap = -20
 
  return tsg_ppt


def fill_slide_not_common(tsg_ppt, org_all, a, c1, c2, c3, c4, c5,data ,box=None):
  fill_slide_title(tsg_ppt, a, org_all)
  not_common_slide=tsg_ppt.slides[a]
  mittelwert = not_common_slide.placeholders[18]
  text_frame = mittelwert.text_frame
  text_frame.clear()
  p = text_frame.paragraphs[0]
  run = p.add_run()
  run.text=box
  font = run.font
  font.size=Pt(10.5)
  font.color.rgb = RGBColor(226, 0, 116)

  chart_data = ChartData()
  chart_data.categories= [c1, c2, c3, c4, c5]
  chart_data.add_series('01', (data[0],data[1],data[2],data[3],data[4]))
  x,y,cx,cy = Inches(0.3), Inches(2.25), Inches(9.38), Inches(4.7)
  graphic_frame = not_common_slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  value_axis.maximum_scale = 100.0
  tick_labels = value_axis.tick_labels
  tick_labels.number_format = '0"%"'
  tick_labels.font.size = Pt(12)
  tick_labels.font.type = 'Tele-GroteskEENor'
  plot = chart.plots[0]
  plot.has_data_labels = True
  data_labels = plot.data_labels
  data_labels.font.size = Pt(12)
  data_labels.number_format = '0"%"'
  data_labels.font.color.rgb = RGBColor(0, 0, 0)
  data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  bar_plot = chart.plots[0]
  bar_plot.gap_width = 20
  bar_plot.overlap = -20
  
  chart.replace_data(chart_data)
  return tsg_ppt

def fill_slide_not_common_2(tsg_ppt, org_all, a, c1, c2, c3, c4, c5,data, box=None):
  fill_slide_title(tsg_ppt, a, org_all)
  not_common_slide=tsg_ppt.slides[a]
  mittelwert = not_common_slide.placeholders[18]
  text_frame = mittelwert.text_frame
  text_frame.clear()
  p = text_frame.paragraphs[0]
  run = p.add_run()
  run.text=box
  font = run.font
  font.size=Pt(10.5)
  font.color.rgb = RGBColor(226, 0, 116)

  chart_data = ChartData()
  chart_data.categories= [c1,c2,c3,c4,c5]
  chart_data.add_series('01', (data[0],data[1],data[2],data[3],data[4]))
  x,y,cx,cy = Inches(0.3), Inches(2.25), Inches(9.38), Inches(4.7)
  graphic_frame = not_common_slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  value_axis.maximum_scale = 100.0
  tick_labels = value_axis.tick_labels
  tick_labels.number_format = '0"%"'
  tick_labels.font.size = Pt(12)
  tick_labels.font.type = 'Tele-GroteskEENor'
  plot = chart.plots[0]
  plot.has_data_labels = True
  data_labels = plot.data_labels
  data_labels.font.size = Pt(12)
  data_labels.number_format = '0"%"'
  data_labels.font.color.rgb = RGBColor(0, 0, 0)
  data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  bar_plot = chart.plots[0]
  bar_plot.gap_width = 20
  bar_plot.overlap = -20
  
  chart.replace_data(chart_data)
  return tsg_ppt
def fill_slide_mean(tsg_ppt, org_all, a, a1, a2, a3, a4, a5, a6, a7, data, box=None):
  common_slide=tsg_ppt.slides[a]
  textbox=common_slide.shapes[0]
  chart_data = ChartData()
  chart_data.categories= [a1,a2,a3,a4,a5,a6,a7]
  chart_data.add_series('01', (data[8], data[7], data[6], data[5], data[4], data[3], data[0]))
  x,y,cx,cy = Inches(0.3), Inches(2.25), Inches(9.38), Inches(4.7)
  graphic_frame = common_slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  value_axis.maximum_scale = 5
  value_axis.minimum_scale = 1
  tick_labels = value_axis.tick_labels
  tick_labels.number_format = '0.00'
  tick_labels.font.size = Pt(12)
  tick_labels.font.type = 'Tele-GroteskEENor'
  plot = chart.plots[0]
  plot.has_data_labels = True
  data_labels = plot.data_labels
  data_labels.font.size = Pt(12)
  data_labels.number_format = '0.00'
  data_labels.font.color.rgb = RGBColor(0, 0, 0)
  data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  bar_plot = chart.plots[0]
  bar_plot.gap_width = 20
  #bar_series = chart.series.bar_series
  #bar_series.fill.solid()
  #bar_series.fill.color.rgb = RGBColor(0, 0, 0)
  

  common_slide=tsg_ppt.slides[a]
  return tsg_ppt

#org is the organization, lista is the list on level org-1, listb is the list with the answers
def getchildresults(org, lista, listb): 
  children = []
  results = []
  for child in structure.orgstructure[org]['child']:
    children.append(child)
  print(children)
  for child in children:
    results.append(listb[lista.index(child)])
    #print(listb.index(child))
  return children, results
    









def id_to_questiontexts(v_id):
  question_ids=['v_175', 'v_176', 'v_177', 'v_179', 'v_180', 'v_181', 'v_182', 'v_183', 'v_184']
  question_texts=["Ich erhalte von meiner Führungskraft regelmäßig Rückmeldungen zu meiner Leistung.", "Wie häufig erhalte ich Rückmeldung zu meiner Leistung von meiner Führungskraft?","Wie lange dauert das Feedbackgespräch durchschnittlich?","Ich erhalte Rückmeldungen zu meiner Gesprächsführung im Kundenkontakt (interner/externer Kunde).", "Ich erhalte Rückmeldungen zu meiner Gesprächsführung im Kundenkontakt (interner/externer Kunde).", "Das Feedback hilft mir, mein Verhalten zu verändern.", "Ich erhalte Feedback zu meinem Beitrag zum Teamerfolg.", "Am Ende des Feedbackgesprächs werden Absprachen getroffen.", "Ich gebe meiner Führungskraft Feedback."]
  for a in range(0,8):
    if (question_ids[a]==v_id):
      return question_texts[a]
  return -1




#fill_slide_common('asdf',1,1,1)

#def asdf(data):
#  mylist = list()
#  for rows in data:
#    if rows[i] not in data and rows[i] and rows[i] not '-99':


##################MAIN###########################

#tsg_ppt=Presentation('tsg_templ_uj_8-16.pptx')

#while (van_graph_element):
#fill_slide_title(tsg_ppt, 2, "Telekom Shop Vertriebsgesellschaft mbH")
##fill_slide_common(data, tsg_ppt, org_all, a, box=None)
mean='2.75'
gultigantwort='600'
#fill_slide_common(1, tsg_ppt, "Telekom Shop Vertriebsgesellschaft mbH", 3, "Trifft überhaupt nicht zu", "Trifft eher nicht zu", "Teils-teils", "trifft eher zu", "Trifft voll zu",  'Mittelwert auf fünfstufiger Skala:'+'\n'+'2,79'+'\n'+'N=300')
org_text_short="TSG"

#--------
if(os.path.isdir('out2')):
  try:
    shutil.rmtree('out2')
  except OSError, e:
    sys.exit("valami fos... in deleting out dir\nError: %s - %s." % (e.filename,e.strerror))
#create out dir, if exists
try:
  os.mkdir('out2')
except OSError, e:
  sys.exit("valami fos... in creating out dir\nError: %s - %s." % (e.filename,e.strerror))



file1_in = open(sys.argv[1], 'rU')
try:
  data_reader = csv.reader((line.replace('\0','') for line in file1_in), delimiter=';')
  my_content = list(data_reader)
  my_header = my_content[0]
  del my_content[0]
finally:
  file1_in.close()


############
#Fill in org lists

level5_users,level5_numbers,level5_filledin_users = structure.fill_in_users(structure.list5, "u_org_5", structure.users_number5)
list_full_5 = (structure.printout(structure.list5,level5_users))
filled_list_5,my_sums_5,my_means_5 = structure.count_percent(list_full_5, level5_numbers)

level4_users,level4_numbers,level4_filledin_users = structure.fill_in_users(structure.list4, "u_org_4", structure.users_number4)
list_full_4 = (structure.printout(structure.list4,level4_users))
filled_list_4,my_sums_4,my_means_4 = structure.count_percent(list_full_4, level4_numbers)

level3_users,level3_numbers,level3_filledin_users = structure.fill_in_users(structure.list3, "u_org_3", structure.users_number3)
list_full_3 = (structure.printout(structure.list3,level3_users))
filled_list_3,my_sums_3,my_means_3 = structure.count_percent(list_full_3, level3_numbers)

level2_users,level2_numbers,level2_filledin_users = structure.fill_in_users(structure.list2, "u_org_2", structure.users_number2)
list_full_2 = (structure.printout(structure.list2,level2_users))
filled_list_2,my_sums_2,my_means_2 = structure.count_percent(list_full_2, level2_numbers)

level1_users,level1_numbers,level1_filledin_users = structure.fill_in_users(structure.list1, "u_org_1", structure.users_number1)
list_full_1 = (structure.printout(structure.list1,level1_users))
filled_list_1,my_sums_1,my_means_1 = structure.count_percent(list_full_1, level1_numbers)
print(level5_filledin_users)

#lista1, lista2 = getchildresults('TSG', structure.list2, filled_list_2)

#structure.fill_in_org(level5_filledin_users, level4_filledin_users, level3_filledin_users, level2_filledin_users, level1_filledin_users)

def fill_in_orgstruct_questions():
  for org in structure.list1:
    for ij, qq in enumerate(['q1', 'q2', 'q3', 'q4', 'q5', 'q6', 'q7', 'q8', 'q9']):
      structure.orgstructure[org][qq] = filled_list_1[structure.list1.index(org)][ij]
    structure.orgstructure[org]['filled_percent'] = round(float(level1_filledin_users[structure.list1.index(org)]) / float(level1_numbers[structure.list1.index(org)])*100)
    structure.orgstructure[org]['filled_in_users'] = level1_filledin_users[structure.list1.index(org)]

  for org in structure.list2:
    for ij, qq in enumerate(['q1', 'q2', 'q3', 'q4', 'q5', 'q6', 'q7', 'q8', 'q9']):
      structure.orgstructure[org][qq] = filled_list_2[structure.list2.index(org)][ij]
    structure.orgstructure[org]['filled_percent'] = round(float(level2_filledin_users[structure.list2.index(org)]) / float(level2_numbers[structure.list2.index(org)])*100)
    structure.orgstructure[org]['filled_in_users'] = level2_filledin_users[structure.list2.index(org)]

  for org in structure.list3:
    for ij, qq in enumerate(['q1', 'q2', 'q3', 'q4', 'q5', 'q6', 'q7', 'q8', 'q9']):
      structure.orgstructure[org][qq] = filled_list_3[structure.list3.index(org)][ij]
    structure.orgstructure[org]['filled_percent'] = round(float(level3_filledin_users[structure.list3.index(org)]) / float(level3_numbers[structure.list3.index(org)])*100)
    structure.orgstructure[org]['filled_in_users'] = level3_filledin_users[structure.list3.index(org)]

  for org in structure.list4:
    for ij, qq in enumerate(['q1', 'q2', 'q3', 'q4', 'q5', 'q6', 'q7', 'q8', 'q9']):
      structure.orgstructure[org][qq] = filled_list_4[structure.list4.index(org)][ij]
    structure.orgstructure[org]['filled_percent'] = round(float(level4_filledin_users[structure.list4.index(org)]) / float(level4_numbers[structure.list4.index(org)])*100)
    structure.orgstructure[org]['filled_in_users'] = level4_filledin_users[structure.list4.index(org)]

  for org in structure.list5:
    for ij, qq in enumerate(['q1', 'q2', 'q3', 'q4', 'q5', 'q6', 'q7', 'q8', 'q9']):
      structure.orgstructure[org][qq] = filled_list_5[structure.list5.index(org)][ij]
    structure.orgstructure[org]['filled_percent'] = round(float(level5_filledin_users[structure.list5.index(org)]) / float(level5_numbers[structure.list5.index(org)])*100)
    structure.orgstructure[org]['filled_in_users'] = level5_filledin_users[structure.list5.index(org)]

#print(structure.orgstructure)

fill_in_orgstruct_questions()

def get_percent(org_list):
  my_percent_list = []
  for orgs in org_list:
    my_percent_list.append(structure.orgstructure[orgs]['filled_percent'])
  return my_percent_list

def fill_slide_1(org, tsg_ppt):
  #print(org)
  orglongname = org['long name'] #orgs.org_long[orgs.org_short.index(org)]
  first_slide = tsg_ppt.slides[0]

  org_1 = first_slide.placeholders[1]
  org_1.text = orglongname + "\n"+(time.strftime("%d.%m.%Y"))

def asdftemp(org_name, tsg_ppt):
  my_level = structure.orgstructure[org_name]['level']
  if my_level == 2: # 2 1 3
    my_orgs = [org_name, structure.orgstructure[org_name]['parent']]
    my_orgs.extend(structure.orgstructure[org_name]['child'])
    #print(my_orgs)

def fill_slide_2(org_name, org, tsg_ppt, my_orgs):
  #own, one above, first, one below
  my_percents = get_percent(my_orgs)
  slide = tsg_ppt.slides[1]
  asdf_text = slide.placeholders[17]
  asdf_text.text = org['long name']
  chart_data = ChartData()
  chart_data.categories= my_orgs #org_names
  series=chart_data.add_series('Rücklaufsquote', my_percents)
  x,y,cx,cy = Cm(dx), Cm(dy), Cm(dcx_2), Cm(dcy)
  graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data)
  chart = graphic_frame.chart
  chart.series[0].fill.solid()
  chart.series[0].fill.fore_color.rgb = RGBColor(82, 154, 214)
  value_axis = chart.value_axis
  value_axis.maximum_scale = 100.0
  plot = chart.plots[0]
  plot.has_data_labels = True
  data_labels = plot.data_labels
  data_labels.font.size = Pt(12)
  data_labels.font.bold = True
  data_labels.font.color.rgb = RGBColor(0, 0, 0)
  data_labels.position = XL_LABEL_POSITION.INSIDE_BASE
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.tick_labels.font.size = Pt(12)
  bar_plot = chart.plots[0]
  #bar_plot.ChartFormat()
  bar_plot.gap_width = g
  bar_plot.overlap = -20
  value_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.major_tick_mark = XL_TICK_MARK.NONE
  value_axis.major_tick_mark = XL_TICK_MARK.NONE
  category_axis.has_major_gridlines = False
  value_axis.has_major_gridlines = False
  value_axis.visible = False
  #bar_plot.fill.solid()
  #bar_plot.fill.fore_color.rgb=RGBColor(0,0,0)

def fill_slide_3(org_name, org, tsg_ppt, my_orgs): #, diff):
  #5 chart
  #1. chart_data: q1, q2 válaszok (1+2,3,4+5) ; a: q4, b: q1
  # c: q6, d: q5
  #q9, q8, q7
  #g: q2
  #h: q3
  #diff: org[0]+org[1] - org_last_yr
  diff = int(round(structure.orgstructure[org_name]['q1'][3] + structure.orgstructure[org_name]['q1'][4]) - round(orgs.org_last_yr[org_name]))
  slide = tsg_ppt.slides[2]
  g_3 = 60
  asdf_text = slide.placeholders[0]
  asdf_text.text = "Leadership Survey@TSG" + "\n" + "GESAMTSICHT " + "\n" + org['long name']
  chart_data = ChartData()
  chart_data_2 =  ChartData()
  chart_data_3 =  ChartData()
  chart_data_4 = ChartData()
  chart_data_5 = ChartData()
  chart_data.categories=[org['long name'], org['long name']]
  chart_data_2.categories=[org['long name'], org['long name']]
  chart_data_3.categories=[org['long name'], org['long name'], org['long name']]
  chart_data_4.categories=[org['long name'], org['long name']]
  chart_data_5.categories=[org['long name'], org['long name']]
  #print(org['q4'][0]+org['q4'][1], org['q1'][0]+org['q1'][1])
  chart_data.add_series('Trifft überhaupt nicht zu / Trifft eher nicht zu', tuple([org['q4'][0]+org['q4'][1], org['q1'][0]+org['q1'][1]])) #(a[0], b[0]))
  chart_data.add_series('Teils-teils', tuple([org['q4'][2], org['q1'][2]])) #(a[1], b[1]))
  chart_data.add_series('Trifft eher zu / Trifft voll zu', tuple([org['q4'][3]+org['q4'][4], org['q1'][3]+org['q1'][4]])) #(a[2], b[2]))
  chart_data_2.add_series('Trifft überhaupt nicht zu / Trifft eher nicht zu', tuple([org['q6'][0]+org['q6'][1], org['q5'][0]+org['q5'][1]])) #(c[0], d[0]))
  chart_data_2.add_series('Teils-teils',tuple([org['q6'][2], org['q5'][2]]))  #(c[1], d[1]))
  chart_data_2.add_series('Trifft eher zu / Trifft voll zu',tuple([org['q6'][3]+org['q6'][4], org['q5'][3]+org['q5'][4]])) #(c[2], d[2]))
  chart_data_3.add_series('Trifft überhaupt nicht zu / Trifft eher nicht zu',tuple([org['q9'][0]+org['q9'][1], org['q8'][0]+org['q8'][1], org['q7'][0]+org['q7'][1]])) #(e[0], f[0], g[0]))
  chart_data_3.add_series('Teils-teils',tuple([org['q9'][2], org['q8'][2], org['q7'][2]])) #(e[1], f[1], g[1]))
  chart_data_3.add_series('Trifft eher zu / Trifft voll zu',tuple([org['q9'][3]+org['q9'][4], org['q8'][3]+org['q8'][4], org['q7'][3]+org['q7'][4]])) #(e[2], f[2], g[2]))
  chart_data_4.add_series('täglich',tuple([0, org['q2'][0]]))
  chart_data_4.add_series('maximal 1x pro Woche',tuple([0, org['q2'][1]]))
  chart_data_4.add_series('bis zu 1x pro Monat',tuple([0, org['q2'][2]]))
  chart_data_4.add_series('halbjährlich',tuple([0, org['q2'][3]]))
  chart_data_4.add_series('seltener',tuple([0, org['q2'][4]]))
  chart_data_5.add_series('1-3 min',tuple([0, org['q3'][0]]))
  chart_data_5.add_series('3-5 min',tuple([0, org['q3'][1]]))
  chart_data_5.add_series('5-15 min',tuple([0, org['q3'][2]]))
  chart_data_5.add_series('15-30 min',tuple([0, org['q3'][3]]))
  chart_data_5.add_series('länger',tuple([0, org['q3'][4]]))
  #chart_data_4.add_series('4',(h[5], g[5]))

  x,y,cx,cy = Cm(dx_3), Cm(dy_3_1), Cm(dcx_3), Cm(dcy_3_1)
  graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  chart.plots[0].has_data_labels = True
  data_labels = chart.plots[0].data_labels
  bar_plot = chart.plots[0]
  bar_plot.gap_width = g_3
  #nehasznaaalddata_labels.position = XL_LABEL_POSITION.OUTSIDE_END
  chart.has_legend = False
  data_labels.font.size = Pt(12)
  data_labels.number_format = nf
  data_labels.font.color.rgb = RGBColor(0, 0, 0)
  data_labels.font.bold = True
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  value_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.major_tick_mark = XL_TICK_MARK.NONE
  value_axis.major_tick_mark = XL_TICK_MARK.NONE
  category_axis.has_major_gridlines = False
  value_axis.has_major_gridlines = False
  value_axis.visible = False
  category_axis.visible = False
  x,y,cx,cy = Cm(dx_3), Cm(dy_3_2), Cm(dcx_3), Cm(dcy_3_1)
  graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data_2)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  chart.plots[0].has_data_labels = True
  data_labels = chart.plots[0].data_labels
  bar_plot = chart.plots[0]
  bar_plot.gap_width = g_3
  data_labels.font.bold = True
  chart.has_legend = False
  data_labels.font.size = Pt(12)
  data_labels.number_format = nf
  data_labels.font.color.rgb = RGBColor(0, 0, 0)
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  value_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.major_tick_mark = XL_TICK_MARK.NONE
  value_axis.major_tick_mark = XL_TICK_MARK.NONE
  category_axis.has_major_gridlines = False
  value_axis.has_major_gridlines = False
  value_axis.visible = False
  category_axis.visible = False

  x,y,cx,cy = Cm(dx_3), Cm(dy_3_3), Cm(dcx_3), Cm(dcy_3_2)
  graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data_3)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  chart.plots[0].has_data_labels = True
  data_labels = chart.plots[0].data_labels
  bar_plot = chart.plots[0]
  bar_plot.gap_width = g_3
  data_labels.font.bold = True
  chart.has_legend = False
  data_labels.font.size = Pt(12)
  data_labels.number_format = nf
  data_labels.font.color.rgb = RGBColor(0, 0, 0)
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  value_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.major_tick_mark = XL_TICK_MARK.NONE
  value_axis.major_tick_mark = XL_TICK_MARK.NONE
  category_axis.has_major_gridlines = False
  value_axis.has_major_gridlines = False
  value_axis.visible = False
  category_axis.visible = False

  x,y,cx,cy = Cm(dx_3), Cm(dy_3_4), Cm(dcx_3), Cm(dcy_3_1)
  graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data_4)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  chart.plots[0].has_data_labels = True
  data_labels = chart.plots[0].data_labels
  data_labels.font.bold = True
  bar_plot = chart.plots[0]
  bar_plot.gap_width = g_3
  #nehasznaaalddata_labels.position = XL_LABEL_POSITION.OUTSIDE_END
  chart.has_legend = False
  data_labels.font.size = Pt(12)
  data_labels.number_format = nf
  data_labels.font.color.rgb = RGBColor(255, 255, 255)
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  value_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.major_tick_mark = XL_TICK_MARK.NONE
  value_axis.major_tick_mark = XL_TICK_MARK.NONE
  category_axis.has_major_gridlines = False
  value_axis.has_major_gridlines = False
  value_axis.visible = False
  category_axis.visible = False
  chart.series[0].fill.solid()
  chart.series[0].fill.fore_color.rgb = RGBColor(8, 49, 87)
  chart.series[1].fill.solid()
  chart.series[1].fill.fore_color.rgb = RGBColor(12, 74, 130)
  chart.series[2].fill.solid()
  chart.series[2].fill.fore_color.rgb = RGBColor(164, 164, 164)
  chart.series[3].fill.solid()
  chart.series[3].fill.fore_color.rgb = RGBColor(149, 159, 44)
  chart.series[4].fill.solid()
  chart.series[4].fill.fore_color.rgb = RGBColor(99, 106, 29)

  x,y,cx,cy = Cm(dx_3), Cm(dy_3_5), Cm(dcx_3), Cm(dcy_3_1)
  graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data_5)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  chart.plots[0].has_data_labels = True
  data_labels = chart.plots[0].data_labels
  data_labels.font.bold = True
  bar_plot = chart.plots[0]
  bar_plot.gap_width = g_3
  #nehasznaaalddata_labels.position = XL_LABEL_POSITION.OUTSIDE_END
  chart.has_legend = False
  data_labels.font.size = Pt(12)
  data_labels.number_format = nf
  data_labels.font.color.rgb = RGBColor(255, 255, 255)
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  value_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.major_tick_mark = XL_TICK_MARK.NONE
  value_axis.major_tick_mark = XL_TICK_MARK.NONE
  category_axis.has_major_gridlines = False
  value_axis.has_major_gridlines = False
  value_axis.visible = False
  category_axis.visible = False
  chart.series[0].fill.solid()
  chart.series[0].fill.fore_color.rgb = RGBColor(14, 87, 81)
  chart.series[1].fill.solid()
  chart.series[1].fill.fore_color.rgb = RGBColor(20, 130, 122)
  chart.series[2].fill.solid()
  chart.series[2].fill.fore_color.rgb = RGBColor(147, 147, 147)
  chart.series[3].fill.solid()
  chart.series[3].fill.fore_color.rgb = RGBColor(222, 176, 0)
  chart.series[4].fill.solid()
  chart.series[4].fill.fore_color.rgb = RGBColor(148, 118, 0)

#chart.replace_data(chart_data)
  
  
  rows=1
  cols=1
  left = Cm(23.15)
  top = Cm(4.52)
  width = Cm(0.66)
  height = Cm(0.66)# set column widths
  table = slide.shapes.add_table(rows, cols, left, top, width, height).table
  table.columns[0].width = Cm(1)
  if diff > 0:
    table.cell(0, 0).text = "+" + str(diff)
  elif diff < 0:
    table.cell(0, 0).text = " " + str(diff)
  elif diff==0:
    table.cell(0, 0).text = "  " + str(diff)
  left = Cm(24.3)
  height = Cm(0.6)
  top = Cm(4.52)
  #pic = slide.shapes.add_picture(img_path, left, top, height=height)
  #return tsg_ppt
  #table.cell(0, 0).text = str(d)
  cell = table.rows[0].cells[0]
  paragraph = cell.textframe.paragraphs[0]
  paragraph.font.size = Pt(be)
  paragraph.font.color.rgb = RGBColor(255, 255, 255)
  #cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
  cell.vertical_anchor = MSO_ANCHOR.MIDDLE
  cell.fill.solid()
  cell.fill.fore_color.rgb = RGBColor(124,124,124)
  if diff > 0:
    img_path='zoldnyil.png'
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
  elif diff < 0:
    img_path='pirosnyil.png'
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
  else:
    top = Cm(4.65)
    img_path='keknyil.png'
    height= Cm(bah)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
  return tsg_ppt

def series_from_orglist_1(orglist, question):
  lista = tuple()
  for my_org in orglist:
    lista = lista + (round(structure.orgstructure[my_org][question][0] + structure.orgstructure[my_org][question][1]), )
  return lista


def series_from_orglist_2(orglist, question):
  lista = tuple()
  for my_org in orglist:
    lista = lista + (round(structure.orgstructure[my_org][question][2]), )
  return lista


def series_from_orglist_3(orglist, question):
  lista = tuple()
  for my_org in orglist:
    lista = lista + (round(structure.orgstructure[my_org][question][3] + structure.orgstructure[my_org][question][4]), )
  return lista


def fill_slide_4(org_name, org, tsg_ppt, my_orgs):
  #1. kérdéshez a válaszok
  #+ diff -- TODO

  my_level = structure.orgstructure[org_name]['level']

  n = len(my_orgs)
  slide = tsg_ppt.slides[3]  
  asdf_text = slide.placeholders[17]
  asdf_text.text = org['long name']
  chart_data = ChartData()
  chart_data.categories = my_orgs
  chart_data.add_series('Trifft überhaupt nicht zu / Trifft eher nicht zu', series_from_orglist_1(my_orgs, 'q1'))
  chart_data.add_series('Teils-teils', series_from_orglist_2(my_orgs, 'q1'))
  chart_data.add_series('Trifft eher zu / Trifft voll zu', series_from_orglist_3(my_orgs, 'q1'))
  
  if len(my_orgs)==3: #n: orgs tömb
    my_bar_plot_gap_width = g
  elif n==5:
    my_bar_plot_gap_width = g
  elif n==8:
    my_bar_plot_gap_width = g
  elif n==9:
    my_bar_plot_gap_width = g
  elif n==11:  
    my_bar_plot_gap_width = g
  elif n==13:  
    my_bar_plot_gap_width = g
  elif n==14:  
    my_bar_plot_gap_width = g
  elif n==15:  
    my_bar_plot_gap_width = g
  elif n==16:  
    my_bar_plot_gap_width = g
  else:
    print('n not in 5,8,11,13,14,15,16')

  x,y,cx,cy = Cm(dx), Cm(dy), Cm(dcx_4), Cm(dcy)
  graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  chart.plots[0].has_data_labels = True
  data_labels = chart.plots[0].data_labels
  bar_plot = chart.plots[0]
  bar_plot.gap_width = my_bar_plot_gap_width
  chart.has_legend = False
  data_labels.font.size = Pt(12)
  data_labels.number_format = nf
  data_labels.font.color.rgb = RGBColor(0, 0, 0)
  data_labels.font.bold = True
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  value_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.major_tick_mark = XL_TICK_MARK.NONE
  value_axis.major_tick_mark = XL_TICK_MARK.NONE
  category_axis.has_major_gridlines = False
  value_axis.has_major_gridlines = False
  value_axis.visible = False
 
 
  #4-5 összevont érték és diff dict-ben levő érték külömbsége kiirva
  #org - tavalyi
  for i,org in enumerate(my_orgs[::-1]):
    #if org == 'TSG':
    #print(round(structure.orgstructure[org]['q1'][3]), round(structure.orgstructure[org]['q1'][4]), round(orgs.org_last_yr[org]))
    diff = int((round(structure.orgstructure[org]['q1'][3] +structure.orgstructure[org]['q1'][4]) - round(orgs.org_last_yr[org])))
    height = Cm(ph)
    diff_text = " "
    if diff > 0:
      img_path='zoldnyil.png'
      diff_text = "+"
    elif diff < 0:
      img_path='pirosnyil.png'
    else:
      img_path='keknyil.png'
      height = Cm(bah)
      diff_text = "  "
    if n==3:
      top_table = Cm(t+i*he)
      top_pic = Cm(t+i*he)
      if diff == 0:
        top_pic = Cm((t+i*he)+0.2)
        #top_table = Cm(t-0.02+i*he)
    elif n==5:
      top_table = Cm(t+i*he)
      top_pic = Cm(t+i*he)
      if diff == 0:
        top_pic = Cm((t+i*he)+0.2)
    elif n==8:
      top_table = Cm(t+i*he)
      top_pic = Cm(t+i*he)
      if diff == 0:
        top_pic = Cm((t+i*he)+0.14)
    elif n==9:
      top_table = Cm(t+i*he)
      top_pic = Cm(t+i*he)
      if diff == 0:
        top_pic = Cm((t+i*he)+0.14)
    elif n==11:
      top_table = Cm(t+i*he)
      top_pic = Cm(t+i*he)
      if diff == 0:
        top_pic = Cm((t+i*he)+0.14)
    elif n==13:
      top_table = Cm(t+i*he)
      top_pic = Cm(t+i*he)
      if diff == 0:
        top_pic = Cm((t+i*he)+0.14)
    elif n==14:
      top_table = Cm(t+i*he)
      top_pic = Cm(t+i*he)
      if diff == 0:
        top_pic = Cm((t+i*he)+0.14)
    elif n==15:
      top_table = Cm(t+i*he)
      top_pic = Cm(t+i*he)
      if diff == 0:
        top_pic = Cm((t+i*he)+0.14)
    elif n==16:
      top_table = Cm(t+i*he)
      top_pic = Cm(t+i*he)
      if diff == 0:
        top_pic = Cm((t+i*he)+0.14)
    left = Cm(aleft_4)
    pic = slide.shapes.add_picture(img_path, left, top_pic,  height=height)
    rows = 1
    cols = 1
    left = Cm(dleft_4)
    width = Cm(1)
    height = Cm(ch)# set column widths
    table = slide.shapes.add_table(rows, cols, left, top_table, width, height).table
    table.columns[0].width = Cm(1)
    table.cell(0, 0).text = diff_text + str(diff)
    cell = table.rows[0].cells[0]
    paragraph = cell.textframe.paragraphs[0]
    paragraph.font.size = Pt(be)
    paragraph.font.color.rgb = RGBColor(255, 255, 255)
    cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(124,124,124)


def fill_slide_5_to_10(org_name, question, s, tsg_ppt, my_orgs): #org_name, q1, 5, tsg_ppt
  #4-5-6-7-8-9 kérdés minden szervezetre (a,b,c) a legutolsó, k: self
  #k0: org['q4'][0]+org['q4'][1]
  #
  slide = tsg_ppt.slides[s-1]
  my_level = structure.orgstructure[org_name]['level']

  asdf_text = slide.placeholders[17]
  asdf_text.text = structure.orgstructure[org_name]['long name']
  chart_data = ChartData()
  chart_data.categories = my_orgs
  my_orgs_answer1 = tuple()
  my_orgs_answer2 = tuple()
  my_orgs_answer3 = tuple()
  for n, org in enumerate(my_orgs):
    my_orgs_answer1 = my_orgs_answer1 + (int(round(structure.orgstructure[org][question][0] + structure.orgstructure[org][question][1])),)
    my_orgs_answer2 = my_orgs_answer2 + (int(round(structure.orgstructure[org][question][2])),)
    my_orgs_answer3 = my_orgs_answer3 + (int(round(structure.orgstructure[org][question][3] + structure.orgstructure[org][question][4])),)
  chart_data.add_series('Trifft überhaupt nicht zu / Trifft eher nicht zu', my_orgs_answer1)
  chart_data.add_series('Teils-teils', my_orgs_answer2)
  chart_data.add_series('Trifft eher zu / Trifft voll zu', my_orgs_answer3)
  x,y,cx,cy = Cm(dx), Cm(dy), Cm(dcx), Cm(dcy)
  graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  chart.plots[0].has_data_labels = True
  data_labels = chart.plots[0].data_labels
  data_labels.font.bold = True
  bar_plot = chart.plots[0]
  bar_plot.gap_width = g
  chart.has_legend = False
  #print(chart_data[1][0])
  data_labels.font.size = Pt(12)
  data_labels.number_format = nf
  data_labels.font.color.rgb = RGBColor(0, 0, 0)
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  value_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.major_tick_mark = XL_TICK_MARK.NONE
  value_axis.major_tick_mark = XL_TICK_MARK.NONE
  category_axis.has_major_gridlines = False
  value_axis.has_major_gridlines = False
  value_axis.visible = False


def fill_slide_11_to_12(org_name, question, s, tsg_ppt, my_orgs):
  #utolsó: 2-3 kérdés, külön 5 válasz, minden org-ra
  slide = tsg_ppt.slides[s-1]
  my_level = structure.orgstructure[org_name]['level']
  if question == 'q2':
    m1 = 'täglich'
    m2 = "maximal 1x pro Woche"
    m3 = "bis zu 1x  pro  Monat"
    m4 = "halbjährlich"
    m5 = "seltener"
  else:
    m1 = '1-3 min'
    m2 = "3-5 min"
    m3 = "5-15 min"
    m4 = "15-30 min"
    m5 = "länger"
 
  my_orgs_answer1 = tuple()
  my_orgs_answer2 = tuple()
  my_orgs_answer3 = tuple()
  my_orgs_answer4 = tuple()
  my_orgs_answer5 = tuple()
  for n, org in enumerate(my_orgs):
    my_orgs_answer1 = my_orgs_answer1 + (int(round(structure.orgstructure[org][question][0])),)
    my_orgs_answer2 = my_orgs_answer2 + (int(round(structure.orgstructure[org][question][1])),)
    my_orgs_answer3 = my_orgs_answer3 + (int(round(structure.orgstructure[org][question][2])),)
    my_orgs_answer4 = my_orgs_answer4 + (int(round(structure.orgstructure[org][question][3])),)
    my_orgs_answer5 = my_orgs_answer5 + (int(round(structure.orgstructure[org][question][4])),)

  asdf_text = slide.placeholders[17]
  asdf_text.text = structure.orgstructure[org_name]['long name']
  chart_data = ChartData()
  chart_data.categories = my_orgs
  chart_data.add_series( m1, my_orgs_answer1)
  chart_data.add_series( m2, my_orgs_answer2)
  chart_data.add_series( m3, my_orgs_answer3)
  chart_data.add_series( m4, my_orgs_answer4)
  chart_data.add_series( m5, my_orgs_answer5)
  x,y,cx,cy = Cm(dx), Cm(dy), Cm(dcx),  Cm(dcy)
  graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  chart.plots[0].has_data_labels = True
  data_labels = chart.plots[0].data_labels
  data_labels.font.bold = True
  bar_plot = chart.plots[0]
  bar_plot.gap_width = g
  chart.has_legend = False
  data_labels.font.size = Pt(12)
  data_labels.number_format = nf
  data_labels.font.color.rgb = RGBColor(255, 255, 255)
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  value_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.major_tick_mark = XL_TICK_MARK.NONE
  value_axis.major_tick_mark = XL_TICK_MARK.NONE
  category_axis.has_major_gridlines = False
  value_axis.has_major_gridlines = False
  value_axis.visible = False


def create_ppt(org_name, org, my_orgs):
  fill_slide_1(org, tsg_ppt)
  fill_slide_2(org_name, org, tsg_ppt, my_orgs)
  fill_slide_3(org_name, org, tsg_ppt, my_orgs)
  fill_slide_4(org_name, org, tsg_ppt, my_orgs)
  for i, question in enumerate(['q4', 'q5', 'q6', 'q7', 'q8', 'q9']):
    fill_slide_5_to_10(org_name, question, i+5, tsg_ppt, my_orgs) #org_name, q1, 5, tsg_ppt
  for i, question in enumerate(['q2', 'q3']):
    fill_slide_11_to_12(org_name, question, i+11, tsg_ppt, my_orgs) #org_name, q1, 5, tsg_ppt
  tsg_ppt.save("out1/"+(org_name.replace(" ", "_")).replace("/","_")+"_TSG_Leadership_Survey"+".pptx")


#create ppts from new structure:

for org in structure.orgstructure.keys():
  if org != 'TSG' and structure.orgstructure[org]['filled_in_users'] > 4:
    my_level = structure.orgstructure[org]['level']
    if my_level == 2: # 2 1 3
      my_orgs = [org, structure.orgstructure[org]['parent']]
      my_orgs.extend(structure.orgstructure[org]['child'])
    elif my_level == 3: # 3 2 1 4
      my_orgs = [org, structure.orgstructure[org]['parent'][0], structure.orgstructure[org]['parent'][1]]
      my_orgs.extend(structure.orgstructure[org]['child'])
    elif my_level == 4: # 4 3 1 5
      my_orgs = [org, structure.orgstructure[org]['parent'][0], structure.orgstructure[org]['parent'][2]]
      my_orgs.extend(structure.orgstructure[org]['child'])
    elif my_level == 5: # 5 4 1
      my_orgs = [org, structure.orgstructure[org]['parent'][0], structure.orgstructure[org]['parent'][3]]
    my_orgs.reverse()

    n = len(my_orgs)
###minden dia
  #number format
    nf = "0;-0; "
    # diagrammok pos hor, pos vert, width
    dx = 0.76
    dy = 5.08
    dcx = 23.2
    # nyilak height
    ahp = 0.6
    ahz = 0.6
    ahk = 0.6
    bah = 0.41
# doboz betumeret
    be = 8
#3. slided
# diagram pos hor
    dx_3 = 15.4
#diagramok pos vert
    dy_3_1 = 3.98  
    dy_3_2 = 6.03
    dy_3_3 = 8.08
    dy_3_4 = 11.13
    dy_3_5 = 12.15
#diagram width
    dcx_3 = 7.9
#diagram height
    dcy_3_1 = 2.78
    dcy_3_2 = 3.81
#gap height
    g_3 = 60
#nyilleft
    aleft_3 = 24.5
#dobozleft
    dleft_3 = 23.15 
#kiveve 2. slide
    dcx_2 = 13.94 
#kiveve 4. slide 
#diagram width
    dcx_4 = 20.65
#nyilleft
    aleft_4 = 23.05
#dobozleft
    dleft_4 = 21.7
    ph = 0.6
    if n==3:
###minden dia   
  # diagram height
      dcy = 4.59
  # diagram gap height
      g = 60
##4. slide 
  #dobozok kozotti tavolsag betweengap
      he = 1.26
      t = 5.69
  #doboz szelesseg
      ch = 0.86
      tsg_ppt=Presentation('tsg_templ_uj_3.pptx')
    elif n==5:
      dcy = 7.45
      g = 60
      he = 1.33
      t = 5.72
      ch = 0.86
      tsg_ppt=Presentation('tsg_templ_uj_5.pptx')
    else:
      dcy = 12
      tsg_ppt=Presentation('tsg_templ_uj_8-16.pptx')
      if n == 8:
        g = 76
        he = 1.40
        t = 5.75
        ch = 0.86
      elif n == 9:
        g = 56
        he = 1.25
        t = 5.68
        ch = 0.86
      elif n == 11:
        g = 56
        he = 1.02
        t = 5.65
        ch = 0.66
      elif n == 13:
        g = 35
        he = 0.865
        t = 5.56
        ch = 0.66
      elif n == 14:
        g = 28
        he = 0.8
        t = 5.55
        ch = 0.66
      elif n == 15:
        g = 21
        he = 0.75
        t = 5.51
        ch = 0.66
      elif n == 16:
        g = 20 
        he = 0.7
        t = 5.51
        ch = 0.66
    create_ppt(org, structure.orgstructure[org], my_orgs)

print "#############################################################TO DOOOOOOOO!!!!!###################################" + "\n" + "1. clear kamudiffertek by VKG 47 S: RVL,S, sajatmaga" + "\n" + "2. ubetuk fajlnevben" +"A:a-to-z" + "K:ellenorzes"

#print(level1_users)
#print(structure.orgstructure)

#for rows in structure.orgstructure.keys():
#  print(rows + ": ", structure.orgstructure[rows])
##
#######
#create ppt for 5th level
#for asdf,my_organi in enumerate(structure.list5):
#  if (level5_filledin_users[asdf] < 5):
#    print 'not created: '+my_organi
#  else:
#    tsg_ppt=Presentation('tsg_templ.pptx')
#    fill_slide_mean(tsg_ppt, my_organi,  2, "Ich gebe meiner Führungskraft Feedback.", "Am Ende des Feedbackgesprächs werden Absprachen getroffen.", "Ich erhalte Feedback zu meinem Beitrag zum Teamerfolg.", "Das Feedback hilft mir, mein Verhalten zu verändern.", "Das Feedbackgespräch baut auf vorherigem Feedback auf.", "Ich erhalte Rückmeldungen zu meiner Gesprächsführung im Kundenkontakt (interner/externer Kunde).", "Wie häufig erhalte ich Rückmeldung zu meiner Leistung von meiner Führungskraft?", my_means_5[asdf])
#    fill_table_slide(tsg_ppt, structure.list5[asdf], str(round(float(level5_filledin_users[asdf] / float(level5_numbers[asdf]))*100,2)) + '%')
#    for i in range(1, 12):
#      fill_slide_title(tsg_ppt, i, org_long[org_short.index(my_organi)]) #org_long_name
#    for i in range(0, 12):
#      if (6 <= i <=12 or i==3):
#        fill_slide_common(tsg_ppt, org_long[org_short.index(my_organi)], i, "Trifft überhaupt nicht zu", "Trifft eher nicht zu", "Teils-teils", "Trifft eher zu", "Trifft voll zu", filled_list_5[asdf][i-3],'Mittelwert auf fünfstufiger Skala:'+'\n'+str(my_means_5[asdf][i-3])+'\n'+"Gültige Antworten:"+'\n'+str(my_sums_5[asdf][i-3])) #org_long_name
#        #will be:
#        #1. tsg_ppt
#        #2. org name
#        #3. number of slide
#        #??? pass other orgs, or search in function??
#        #fill_slide_common(tsg_ppt, org_long[org_short.index(my_organi)], i)
#      elif (i==4):
#        fill_slide_not_common(tsg_ppt, org_long[org_short.index(my_organi)], i, "täglich", "maximal 1x pro Woche", "bis zu 1x pro Monat", "halbjährlich", "seltener", filled_list_5[asdf][i-3],"\n"+"Gültige Antworten:"+"\n"+str(my_sums_5[asdf][i-3]))
#      elif (i==5):
#        fill_slide_not_common_2(tsg_ppt, org_long[org_short.index(my_organi)], i, "1-3 min", "3-5 min", "5-15 min", "15-30 min", "länger", filled_list_5[asdf][i-3],"\n"+"Gültige Antworten:"+"\n"+str(my_sums_5[asdf][i-3]))
#    tsg_ppt.save("out1/"+(structure.list5[asdf].replace(" ", "_")).replace("/","_")+"_TSG_Leadership_Survey"+".pptx")
#
#######
##create ppt for 4th level
#for asdf,my_organi in enumerate(structure.list4):
#  if (level4_filledin_users[asdf] < 5):
#    print 'not created: '+my_organi
#  else:
#    tsg_ppt=Presentation('tsg_templ.pptx')
#    fill_slide_mean(tsg_ppt, my_organi,  2, "Ich gebe meiner Führungskraft Feedback.", "Am Ende des Feedbackgesprächs werden Absprachen getroffen.", "Ich erhalte Feedback zu meinem Beitrag zum Teamerfolg.", "Das Feedback hilft mir, mein Verhalten zu verändern.", "Das Feedbackgespräch baut auf vorherigem Feedback auf.", "Ich erhalte Rückmeldungen zu meiner Gesprächsführung im Kundenkontakt (interner/externer Kunde).", "Wie häufig erhalte ich Rückmeldung zu meiner Leistung von meiner Führungskraft?",my_means_4[asdf])
#    fill_table_slide(tsg_ppt, level4_numbers[asdf], level4_filledin_users[asdf], str(round(float(level4_filledin_users[asdf] / float(level4_numbers[asdf]))*100,2)) + '%')
#    for i in range(1, 12):
#      fill_slide_title(tsg_ppt, i, org_long[org_short.index(my_organi)]) #org_long_name
#    for i in range(0, 12):
#      if (6 <= i <=12 or i==3):
#        fill_slide_common(tsg_ppt, org_long[org_short.index(my_organi)], i, "Trifft überhaupt nicht zu", "Trifft eher nicht zu", "Teils-teils", "Trifft eher zu", "Trifft voll zu", filled_list_4[asdf][i-3],'Mittelwert auf fünfstufiger Skala:'+'\n'+str(my_means_4[asdf][i-3])+'\n'+"Gültige Antworten:"+'\n'+str(my_sums_4[asdf][i-3])) #org_long_name
#      elif (i==4):
#        fill_slide_not_common(tsg_ppt, org_long[org_short.index(my_organi)], i, "täglich", "maximal 1x pro Woche", "bis zu 1x pro Monat", "halbjährlich", "seltener", filled_list_4[asdf][i-3],"\n"+"Gültige Antworten:"+"\n"+str(my_sums_4[asdf][i-3]))
#      elif (i==5):
#        fill_slide_not_common_2(tsg_ppt, org_long[org_short.index(my_organi)], i, "1-3 min", "3-5 min", "5-15 min", "15-30 min", "länger", filled_list_4[asdf][i-3],"\n"+"Gültige Antworten:"+"\n"+str(my_sums_4[asdf][i-3]))
#    tsg_ppt.save("out1/"+(structure.list4[asdf].replace(" ", "_")).replace("/","_")+"_TSG_Leadership_Survey"+".pptx")
#
#######
##create ppt for 3th level
#for asdf,my_organi in enumerate(structure.list3):
#  if (level3_filledin_users[asdf] < 5):
#    print 'not created: '+my_organi
#  else: 
#    tsg_ppt=Presentation('tsg_templ.pptx')
#    fill_slide_mean(tsg_ppt, my_organi,  2, "Ich gebe meiner Führungskraft Feedback.", "Am Ende des Feedbackgesprächs werden Absprachen getroffen.", "Ich erhalte Feedback zu meinem Beitrag zum Teamerfolg.", "Das Feedback hilft mir, mein Verhalten zu verändern.", "Das Feedbackgespräch baut auf vorherigem Feedback auf.", "Ich erhalte Rückmeldungen zu meiner Gesprächsführung im Kundenkontakt (interner/externer Kunde).", "Wie häufig erhalte ich Rückmeldung zu meiner Leistung von meiner Führungskraft?",my_means_3[asdf])
#    fill_table_slide(tsg_ppt, level3_numbers[asdf], level3_filledin_users[asdf], str(round(float(level3_filledin_users[asdf] / float(level3_numbers[asdf]))*100,2)) + '%')
#    for i in range(1, 12):
#      fill_slide_title(tsg_ppt, i, org_long[org_short.index(my_organi)]) #org_long_name
#    for i in range(0, 12):
#      if (6 <= i <=12 or i==3):
#        fill_slide_common(tsg_ppt, org_long[org_short.index(my_organi)], i, "Trifft überhaupt nicht zu", "Trifft eher nicht zu", "Teils-teils", "Trifft eher zu", "Trifft voll zu", filled_list_3[asdf][i-3],'Mittelwert auf fünfstufiger Skala:'+'\n'+str(my_means_3[asdf][i-3])+'\n'+"Gültige Antworten:"+'\n'+str(my_sums_3[asdf][i-3])) #org_long_name
#      elif (i==4):
#        fill_slide_not_common(tsg_ppt, org_long[org_short.index(my_organi)], i, "täglich", "maximal 1x pro Woche", "bis zu 1x pro Monat", "halbjährlich", "seltener", filled_list_3[asdf][i-3],"\n"+"Gültige Antworten:"+"\n"+str(my_sums_3[asdf][i-3]))
#      elif (i==5):
#        fill_slide_not_common_2(tsg_ppt, org_long[org_short.index(my_organi)], i, "1-3 min", "3-5 min", "5-15 min", "15-30 min", "länger", filled_list_3[asdf][i-3],"\n"+"Gültige Antworten:"+"\n"+str(my_sums_3[asdf][i-3]))
#    tsg_ppt.save("out1/"+(structure.list3[asdf].replace(" ", "_")).replace("/","_")+"_TSG_Leadership_Survey"+".pptx")
#
#######
##create ppt for 2th level
#for asdf,my_organi in enumerate(structure.list2):
#  if (level2_filledin_users[asdf] < 5):
#    print 'not created: '+my_organi
#  else:
#    tsg_ppt=Presentation('tsg_templ.pptx')
#    fill_slide_mean(tsg_ppt, my_organi,  2, "Ich gebe meiner Führungskraft Feedback.", "Am Ende des Feedbackgesprächs werden Absprachen getroffen.", "Ich erhalte Feedback zu meinem Beitrag zum Teamerfolg.", "Das Feedback hilft mir, mein Verhalten zu verändern.", "Das Feedbackgespräch baut auf vorherigem Feedback auf.", "Ich erhalte Rückmeldungen zu meiner Gesprächsführung im Kundenkontakt (interner/externer Kunde).", "Wie häufig erhalte ich Rückmeldung zu meiner Leistung von meiner Führungskraft?",my_means_2[asdf])
#    fill_table_slide(tsg_ppt, level2_numbers[asdf], level2_filledin_users[asdf], str(round(float(level2_filledin_users[asdf] / float(level2_numbers[asdf]))*100,2)) + '%')
#    for i in range(1, 12):
#      fill_slide_title(tsg_ppt, i, org_long[org_short.index(my_organi)]) #org_long_name
#    for i in range(0, 12):
#      if (6 <= i <=12 or i==3):
#        fill_slide_common(tsg_ppt, org_long[org_short.index(my_organi)], i, "Trifft überhaupt nicht zu", "Trifft eher nicht zu", "Teils-teils", "Trifft eher zu", "Trifft voll zu", filled_list_2[asdf][i-3],'Mittelwert auf fünfstufiger Skala:'+'\n'+str(my_means_2[asdf][i-3])+'\n'+"Gültige Antworten:"+'\n'+str(my_sums_2[asdf][i-3])) #org_long_name
#      elif (i==4):
#        fill_slide_not_common(tsg_ppt, org_long[org_short.index(my_organi)], i, "täglich", "maximal 1x pro Woche", "bis zu 1x pro Monat", "halbjährlich", "seltener", filled_list_2[asdf][i-3],"\n"+"Gültige Antworten:"+"\n"+str(my_sums_2[asdf][i-3]))
#      elif (i==5):
#        fill_slide_not_common_2(tsg_ppt, org_long[org_short.index(my_organi)], i, "1-3 min", "3-5 min", "5-15 min", "15-30 min", "länger", filled_list_2[asdf][i-3],"\n"+"Gültige Antworten:"+"\n"+str(my_sums_2[asdf][i-3]))
#    tsg_ppt.save("out1/"+(structure.list2[asdf].replace(" ", "_")).replace("/","_")+"_TSG_Leadership_Survey"+".pptx")
#
######
##create ppt for 1th level
#for asdf,my_organi in enumerate(structure.list1):
#  if (level1_filledin_users[asdf] < 5):
#    print 'not created: '+my_organi
#  else:
#    tsg_ppt=Presentation('tsg_templ.pptx')
#    fill_slide_mean(tsg_ppt, my_organi,  2, "Ich gebe meiner Führungskraft Feedback.", "Am Ende des Feedbackgesprächs werden Absprachen getroffen.", "Ich erhalte Feedback zu meinem Beitrag zum Teamerfolg.", "Das Feedback hilft mir, mein Verhalten zu verändern.", "Das Feedbackgespräch baut auf vorherigem Feedback auf.", "Ich erhalte Rückmeldungen zu meiner Gesprächsführung im Kundenkontakt (interner/externer Kunde).", "Wie häufig erhalte ich Rückmeldung zu meiner Leistung von meiner Führungskraft?",my_means_1[asdf])
#    fill_table_slide(tsg_ppt, level1_numbers[asdf], level1_filledin_users[asdf], str(round(float(level1_filledin_users[asdf] / float(level1_numbers[asdf]))*100,2)) + '%')
#    #fill_slide_not_common(1, tsg_ppt, "Telekom Shop Vertriebgesellschaft mbH", 4, "täglich", "maximal 1x pro Woche", "bis zu 1x pro Monat", "halbjährlich", "seltener", "\n"+"\n"+"N=289")
#    for i in range(1, 12):
#      fill_slide_title(tsg_ppt, i, org_long[org_short.index(my_organi)]) #org_long_name
#    for i in range(0, 12):
#      if (6 <= i <=12 or i==3):
#        fill_slide_common(tsg_ppt, org_long[org_short.index(my_organi)], i, "Trifft überhaupt nicht zu", "Trifft eher nicht zu", "Teils-teils", "Trifft eher zu", "Trifft voll zu", filled_list_1[asdf][i-3],'Mittelwert auf fünfstufiger Skala:'+'\n'+str(my_means_1[asdf][i-3])+'\n'+"Gültige Antworten:"+'\n'+str(my_sums_1[asdf][i-3])) #org_long_name
#      elif (i==4):
#        fill_slide_not_common(tsg_ppt, org_long[org_short.index(my_organi)], i, "täglich", "maximal 1x pro Woche", "bis zu 1x pro Monat", "halbjährlich", "seltener", filled_list_1[asdf][i-3],"\n"+"Gültige Antworten:"+"\n"+str(my_sums_1[asdf][i-3]))
#      elif (i==5):
#        fill_slide_not_common_2(tsg_ppt, org_long[org_short.index(my_organi)], i, "1-3 min", "3-5 min", "5-15 min", "15-30 min", "länger", filled_list_1[asdf][i-3],"\n"+"Gültige Antworten:"+"\n"+str(my_sums_1[asdf][i-3]))
#    tsg_ppt.save("out1/"+(structure.list1[asdf].replace(" ", "_")).replace("/","_")+"_TSG_Leadership_Survey"+".pptx")
#
#################MAIN ENDS HERE###################
#dict to fill in from old file in structure.py
#dict_1 = {'ORG1': 73, 'ORG2': 54}
#dict_2 = {'ORG3': 34, 'ORG34': 234}
