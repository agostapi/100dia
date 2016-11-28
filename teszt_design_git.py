#-*- coding: utf-8 -*-

###############TODO##############
#org - short + long version
#5 szabály tesztelés, kiiratni kieso orgokat level3on lesz ilyen
#test-4levelppts
###############/TODO#############
 
from __future__ import absolute_import, print_function, unicode_literals
import codecs
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_TICK_LABEL_POSITION
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.line import LineFormat
from pptx.dml.chtfmt import ChartFormat
from pptx.dml.fill import FillFormat
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.chart import XL_LEGEND_POSITION
import time
#import structure
import csv
import sys
import os
import shutil
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
reload(sys)  
sys.setdefaultencoding('utf8')


#UTF8Writer = codecs.getwriter('utf8')


#fill the titles in all slide withe the long name of the orgs, add datum to the "0" slide

def fill_slide_0_and_titles(tsg_ppt, orglongname):
  for i in range(0,12):
    if i==0:
      first_slide = tsg_ppt.slides[0]
      org_1 = first_slide.placeholders[1]
      org_1.text = orglongname + "\n"+(time.strftime("%d.%m.%Y"))
    else:
      other_slide = tsg_ppt.slides[i]
      org = other_slide.placeholders[17]
      org.text = orglongname
  return tsg_ppt


# fill the first slide with the percents of the participants pro org

def fill_slide_1(tsg_ppt, org_names, org_participants):
  slide = tsg_ppt.slides[1]
  chart_data = ChartData()
  chart_data.categories= org_names
  series=chart_data.add_series('01', org_participants)
  x,y,cx,cy = Inches(0.3), Inches(2.0), Cm(13.7), Cm(11.9)
  graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data)
  chart = graphic_frame.chart
  chart.series[0].fill.solid()
  chart.series[0].fill.fore_color.rgb = RGBColor(82, 154, 214)
  value_axis = chart.value_axis
  value_axis.maximum_scale = 100.0
  plot = chart.plots[0]
  plot.has_data_labels = True
  data_labels = plot.data_labels
  data_labels.font.size = Pt(12)
  data_labels.font.bold = True
  data_labels.font.color.rgb = RGBColor(0, 0, 0)
  data_labels.position = XL_LABEL_POSITION.INSIDE_BASE
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.tick_labels.font.size = Pt(12)
  bar_plot = chart.plots[0]
  #bar_plot.ChartFormat()
  bar_plot.gap_width = 20
  bar_plot.overlap = -20
  value_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.major_tick_mark = XL_TICK_MARK.NONE
  value_axis.major_tick_mark = XL_TICK_MARK.NONE
  category_axis.has_major_gridlines = False
  value_axis.has_major_gridlines = False
  value_axis.visible = False
  #bar_plot.fill.solid()
  #bar_plot.fill.fore_color.rgb=RGBColor(0,0,0)
  return tsg_ppt
#fills the second slide with 9 charts of the 9 questions, appears a table by the first chart with the diffrence to the last years data and with an arrow which signs the way of the difference
def fill_slide_2(tsg_ppt, diff):
  slide = tsg_ppt.slides[2]
  chart_data = ChartData()
  chart_data_2 =  ChartData()
  chart_data_3 =  ChartData()
  chart_data_4 = ChartData()
  chart_data_5 = ChartData()
  chart_data.categories=['1', '2']
  chart_data_2.categories=['1', '2']
  chart_data_3.categories=['1', '2', '3']
  chart_data_4.categories=['1', '2']
  chart_data_5.categories=['1', '2']
  chart_data.add_series('2',(a[0], b[0]))#, 0, 0))#, b[0], c[0], d[0], e[0], f[0], g[0])) #, d[0]))
  chart_data.add_series('3',(a[1], b[1]))#, 0, 0))#, b[1], c[1], d[1], e[1], f[1], g[1]))
  chart_data.add_series('4',(a[2], b[2]))#, 0, 0))#, b[2], c[2], d[2], e[2], f[2], g[2]))#, d[2]))
  chart_data_2.add_series('2',(c[0], d[0]))
  chart_data_2.add_series('3',(c[1], d[1]))
  chart_data_2.add_series('4',(c[2], d[2]))
  chart_data_3.add_series('2',(e[0], f[0], g[0]))
  chart_data_3.add_series('2',(e[1], f[1], g[1]))
  chart_data_3.add_series('2',(e[2], f[2], g[2]))
  chart_data_4.add_series('3',(0, g[0]))
  chart_data_4.add_series('4',(0, g[1]))
  chart_data_4.add_series('5',(0, g[2]))
  chart_data_4.add_series('6',(0, g[3]))
  chart_data_4.add_series('4',(0, g[4]))
  chart_data_5.add_series('3',(0, h[0]))
  chart_data_5.add_series('4',(0, h[1]))
  chart_data_5.add_series('5',(0, h[2]))
  chart_data_5.add_series('6',(0, h[3]))
  chart_data_5.add_series('4',(0, h[4]))
  #chart_data_4.add_series('4',(h[5], g[5]))

  x,y,cx,cy = Cm(10.5), Cm(4.5), Cm(4.6), Cm(3)
  graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  chart.plots[0].has_data_labels = True
  data_labels = chart.plots[0].data_labels
  bar_plot = chart.plots[0]
  bar_plot.gap_width = 10
  #nehasznaaalddata_labels.position = XL_LABEL_POSITION.OUTSIDE_END
  chart.has_legend = False
  data_labels.font.size = Pt(12)
  data_labels.number_format = '0'
  data_labels.font.color.rgb = RGBColor(0, 0, 0)
  data_labels.font.bold = True
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  value_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.major_tick_mark = XL_TICK_MARK.NONE
  value_axis.major_tick_mark = XL_TICK_MARK.NONE
  category_axis.has_major_gridlines = False
  value_axis.has_major_gridlines = False
  value_axis.visible = False
  category_axis.visible = False
  x,y,cx,cy = Cm(10.5), Cm(6.68), Cm(4.6), Cm(3)
  graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data_2)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  chart.plots[0].has_data_labels = True
  data_labels = chart.plots[0].data_labels
  bar_plot = chart.plots[0]
  bar_plot.gap_width = 10
  data_labels.font.bold = True
  chart.has_legend = False
  data_labels.font.size = Pt(12)
  data_labels.number_format = '0'
  data_labels.font.color.rgb = RGBColor(0, 0, 0)
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  value_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.major_tick_mark = XL_TICK_MARK.NONE
  value_axis.major_tick_mark = XL_TICK_MARK.NONE
  category_axis.has_major_gridlines = False
  value_axis.has_major_gridlines = False
  value_axis.visible = False
  category_axis.visible = False

  x,y,cx,cy = Cm(10.5), Cm(8.86), Cm(4.6), Cm(4.1)
  graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data_3)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  chart.plots[0].has_data_labels = True
  data_labels = chart.plots[0].data_labels
  bar_plot = chart.plots[0]
  bar_plot.gap_width = 10
  data_labels.font.bold = True
  chart.has_legend = False
  data_labels.font.size = Pt(12)
  data_labels.number_format = '0'
  data_labels.font.color.rgb = RGBColor(0, 0, 0)
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  value_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.major_tick_mark = XL_TICK_MARK.NONE
  value_axis.major_tick_mark = XL_TICK_MARK.NONE
  category_axis.has_major_gridlines = False
  value_axis.has_major_gridlines = False
  value_axis.visible = False
  category_axis.visible = False

  x,y,cx,cy = Cm(10.5), Cm(12.18), Cm(4.6), Cm(3)
  graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data_4)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  chart.plots[0].has_data_labels = True
  data_labels = chart.plots[0].data_labels
  data_labels.font.bold = True
  bar_plot = chart.plots[0]
  bar_plot.gap_width = 10
  #nehasznaaalddata_labels.position = XL_LABEL_POSITION.OUTSIDE_END
  chart.has_legend = False
  data_labels.font.size = Pt(12)
  data_labels.number_format = '0'
  data_labels.font.color.rgb = RGBColor(255, 255, 255)
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  value_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.major_tick_mark = XL_TICK_MARK.NONE
  value_axis.major_tick_mark = XL_TICK_MARK.NONE
  category_axis.has_major_gridlines = False
  value_axis.has_major_gridlines = False
  value_axis.visible = False
  category_axis.visible = False
  chart.series[0].fill.solid()
  chart.series[0].fill.fore_color.rgb = RGBColor(8, 49, 87)
  chart.series[1].fill.solid()
  chart.series[1].fill.fore_color.rgb = RGBColor(12, 74, 130)
  chart.series[2].fill.solid()
  chart.series[2].fill.fore_color.rgb = RGBColor(164, 164, 164)
  chart.series[3].fill.solid()
  chart.series[3].fill.fore_color.rgb = RGBColor(149, 159, 44)
  chart.series[4].fill.solid()
  chart.series[4].fill.fore_color.rgb = RGBColor(99, 106, 29)

  x,y,cx,cy = Cm(10.5), Cm(13.28), Cm(4.6), Cm(3)
  graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data_5)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  chart.plots[0].has_data_labels = True
  data_labels = chart.plots[0].data_labels
  data_labels.font.bold = True
  bar_plot = chart.plots[0]
  bar_plot.gap_width = 10
  #nehasznaaalddata_labels.position = XL_LABEL_POSITION.OUTSIDE_END
  chart.has_legend = False
  data_labels.font.size = Pt(12)
  data_labels.number_format = '0'
  data_labels.font.color.rgb = RGBColor(255, 255, 255)
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  value_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.major_tick_mark = XL_TICK_MARK.NONE
  value_axis.major_tick_mark = XL_TICK_MARK.NONE
  category_axis.has_major_gridlines = False
  value_axis.has_major_gridlines = False
  value_axis.visible = False
  category_axis.visible = False
  chart.series[0].fill.solid()
  chart.series[0].fill.fore_color.rgb = RGBColor(14, 87, 81)
  chart.series[1].fill.solid()
  chart.series[1].fill.fore_color.rgb = RGBColor(20, 130, 122)
  chart.series[2].fill.solid()
  chart.series[2].fill.fore_color.rgb = RGBColor(147, 147, 147)
  chart.series[3].fill.solid()
  chart.series[3].fill.fore_color.rgb = RGBColor(222, 176, 0)
  chart.series[4].fill.solid()
  chart.series[4].fill.fore_color.rgb = RGBColor(148, 118, 0)

#chart.replace_data(chart_data)
  
  
  rows=1
  cols=1
  left = Inches(5.83)
  top = Inches(1.94)
  width = Inches(0.8)
  height = Inches (0.425)# set column widths
  table = slide.shapes.add_table(rows, cols, left, top, width, height).table
  table.columns[0].width = Inches(0.39)
  if diff[0] > 0:
    table.cell(0, 0).text = "+" + str(diff[0])
  elif diff[0] < 0:
    table.cell(0, 0).text = " " + str(diff[0])
  elif diff[0]==0:
    table.cell(0, 0).text = " " + str(diff[0])
  left = Inches(6.25)
  height = Inches(0.3)
  top = Inches(2)
  #pic = slide.shapes.add_picture(img_path, left, top, height=height)
  #return tsg_ppt
  #table.cell(0, 0).text = str(d)
  cell = table.rows[0].cells[0]
  paragraph = cell.textframe.paragraphs[0]
  paragraph.font.size = Pt(12)
  paragraph.font.color.rgb = RGBColor(255, 255, 255)
  cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
  cell.vertical_anchor = MSO_ANCHOR.MIDDLE
  cell.fill.solid()
  cell.fill.fore_color.rgb = RGBColor(124,124,124)
  if diff[0] > 0:
    img_path='zoldnyil.png'
    left = Inches(6.25)
    height = Cm(0.6)
    top = Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
  elif diff[0] < 0:
    img_path='pirosnyil.png'
    left = Inches(6.25)
    height = Cm(0.6)
    top = Inches(2)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
  else:
    img_path='keknyil.png'
    left = Inches(6.25)
    height = Cm(0.43)
    top = Inches(2.05)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
  return tsg_ppt

def fill_slide_3(tsg_ppt, n):
  if n==3:
    slide = tsg_ppt.slides[3]  
    chart_data = ChartData()
    chart_data.categories = ['1', '2', '3']
    chart_data.add_series('1', (a[0], b[0], c[0]))
    chart_data.add_series('2', (a[1], b[1], c[1]))
    chart_data.add_series('3', (a[2], b[2], c[2]))
    x,y,cx,cy = Inches(0.3), Inches(2.25), Cm(14.14), Cm(11.39)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    value_axis = chart.value_axis
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    bar_plot = chart.plots[0]
    bar_plot.gap_width = 248
    chart.has_legend = False
    data_labels.font.size = Pt(12)
    data_labels.number_format = '0'
    data_labels.font.color.rgb = RGBColor(0, 0, 0)
    data_labels.font.bold = True
    category_axis = chart.category_axis
    category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    category_axis.tick_labels.font.size = Pt(12)
    value_axis.minor_tick_mark = XL_TICK_MARK.NONE
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    value_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.has_major_gridlines = False
    value_axis.has_major_gridlines = False
    value_axis.visible = False
  elif n==5:  
    slide = tsg_ppt.slides[3]  
    chart_data = ChartData()
    chart_data.categories = ['1', '2', '3']
    chart_data.add_series('1', (a[0], b[0], c[0], d[0], e[0]))
    chart_data.add_series('2', (a[1], b[1], c[1], d[1], e[1]))
    chart_data.add_series('3', (a[2], b[2], c[2], d[2], e[2]))
    x,y,cx,cy = Inches(0.3), Inches(2.25), Cm(14.14), Cm(11.39)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    value_axis = chart.value_axis
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    bar_plot = chart.plots[0]
    bar_plot.gap_width = 108
    chart.has_legend = False
    data_labels.font.size = Pt(12)
    data_labels.number_format = '0'
    data_labels.font.color.rgb = RGBColor(0, 0, 0)
    data_labels.font.bold = True
    category_axis = chart.category_axis
    category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    category_axis.tick_labels.font.size = Pt(12)
    value_axis.minor_tick_mark = XL_TICK_MARK.NONE
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    value_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.has_major_gridlines = False
    value_axis.has_major_gridlines = False
    value_axis.visible = False
  elif n==8:  
    slide = tsg_ppt.slides[3]  
    chart_data = ChartData()
    chart_data.categories = ['1', '2', '3']
    chart_data.add_series('1', (a[0], b[0], c[0], d[0], e[0], f[0], ii[0], ji[0]))
    chart_data.add_series('2', (a[1], b[1], c[1], d[1], e[1], f[1], ii[1], ji[1]))
    chart_data.add_series('3', (a[2], b[2], c[2], d[2], e[2], f[2], ii[2], ji[2]))
    x,y,cx,cy = Inches(0.3), Inches(2.25), Cm(14.14), Cm(11.39)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    value_axis = chart.value_axis
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.font.bold = True
    bar_plot = chart.plots[0]
    bar_plot.gap_width = 33
    chart.has_legend = False
    data_labels.font.size = Pt(12)
    data_labels.number_format = '0'
    data_labels.font.color.rgb = RGBColor(0, 0, 0)
    category_axis = chart.category_axis
    category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    category_axis.tick_labels.font.size = Pt(12)
    value_axis.minor_tick_mark = XL_TICK_MARK.NONE
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    value_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.has_major_gridlines = False
    value_axis.has_major_gridlines = False
    value_axis.visible = False
  elif n==9:  
    slide = tsg_ppt.slides[3]  
    chart_data = ChartData()
    chart_data.categories = ['1', '2', '3']
    chart_data.add_series('1', (a[0], b[0], c[0], d[0], e[0], f[0], ii[0], ji[0], k[0]))
    chart_data.add_series('2', (a[1], b[1], c[1], d[1], e[1], f[1], ii[1], ji[1], k[1]))
    chart_data.add_series('3', (a[2], b[2], c[2], d[2], e[2], f[2], ii[2], ji[2], k[2]))
    x,y,cx,cy = Inches(0.3), Inches(2.25), Cm(14.14), Cm(11.39)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    value_axis = chart.value_axis
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.font.bold = True
    bar_plot = chart.plots[0]
    bar_plot.gap_width = 23
    chart.has_legend = False
    data_labels.font.size = Pt(12)
    data_labels.number_format = '0'
    data_labels.font.color.rgb = RGBColor(0, 0, 0)
    category_axis = chart.category_axis
    category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    category_axis.tick_labels.font.size = Pt(12)
    value_axis.minor_tick_mark = XL_TICK_MARK.NONE
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    value_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.has_major_gridlines = False
    value_axis.has_major_gridlines = False
    value_axis.visible = False
  elif n==11:  
    slide = tsg_ppt.slides[3]  
    chart_data = ChartData()
    chart_data.categories = ['1', '2', '3']
    chart_data.add_series('1', (a[0], b[0], c[0], d[0], e[0], f[0], ii[0], ji[0], k[0], l[0], m[0]))
    chart_data.add_series('2', (a[1], b[1], c[1], d[1], e[1], f[1], ii[1], ji[1], k[1], l[1], m[1]))
    chart_data.add_series('3', (a[2], b[2], c[2], d[2], e[2], f[2], ii[2], ji[2], k[2], l[2], m[2]))
    x,y,cx,cy = Inches(0.3), Inches(2.25), Cm(14.14), Cm(11.39)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    value_axis = chart.value_axis
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.font.bold = True
    bar_plot = chart.plots[0]
    bar_plot.gap_width = 13
    chart.has_legend = False
    data_labels.font.size = Pt(12)
    data_labels.number_format = '0'
    data_labels.font.color.rgb = RGBColor(0, 0, 0)
    category_axis = chart.category_axis
    category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    category_axis.tick_labels.font.size = Pt(12)
    value_axis.minor_tick_mark = XL_TICK_MARK.NONE
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    value_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.has_major_gridlines = False
    value_axis.has_major_gridlines = False
    value_axis.visible = False
  elif n==13:  
    slide = tsg_ppt.slides[3]  
    chart_data = ChartData()
    chart_data.categories = ['1', '2', '3']
    chart_data.add_series('1', (a[0], b[0], c[0], d[0], e[0], f[0], ii[0], ji[0], k[0], l[0], m[0], nn[0], o[0]))
    chart_data.add_series('2', (a[1], b[1], c[1], d[1], e[1], f[1], ii[1], ji[1], k[1], l[1], m[1], nn[1], o[1]))
    chart_data.add_series('3', (a[2], b[2], c[2], d[2], e[2], f[2], ii[2], ji[2], k[2], l[2], m[2], nn[2], o[2]))
    x,y,cx,cy = Inches(0.3), Inches(2.25), Cm(14.14), Cm(11.39)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    value_axis = chart.value_axis
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.font.bold = True
    bar_plot = chart.plots[0]
    bar_plot.gap_width = 13
    chart.has_legend = False
    data_labels.font.size = Pt(12)
    data_labels.number_format = '0'
    data_labels.font.color.rgb = RGBColor(0, 0, 0)
    category_axis = chart.category_axis
    category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    category_axis.tick_labels.font.size = Pt(12)
    value_axis.minor_tick_mark = XL_TICK_MARK.NONE
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    value_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.has_major_gridlines = False
    value_axis.has_major_gridlines = False
    value_axis.visible = False
  elif n==14:  
    slide = tsg_ppt.slides[3]  
    chart_data = ChartData()
    chart_data.categories = ['1', '2', '3']
    chart_data.add_series('1', (a[0], b[0], c[0], d[0], e[0], f[0], ii[0], ji[0], k[0], l[0], m[0], nn[0], o[0], p[0]))
    chart_data.add_series('2', (a[1], b[1], c[1], d[1], e[1], f[1], ii[1], ji[1], k[1], l[1], m[1], nn[1], o[1], p[1]))
    chart_data.add_series('3', (a[2], b[2], c[2], d[2], e[2], f[2], ii[2], ji[2], k[2], l[2], m[2], nn[2], o[2], p[2]))
    x,y,cx,cy = Inches(0.3), Inches(2.25), Cm(14.14), Cm(11.39)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    value_axis = chart.value_axis
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.font.bold = True
    bar_plot = chart.plots[0]
    bar_plot.gap_width = 13
    chart.has_legend = False
    data_labels.font.size = Pt(12)
    data_labels.number_format = '0'
    data_labels.font.color.rgb = RGBColor(0, 0, 0)
    category_axis = chart.category_axis
    category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    category_axis.tick_labels.font.size = Pt(12)
    value_axis.minor_tick_mark = XL_TICK_MARK.NONE
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    value_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.has_major_gridlines = False
    value_axis.has_major_gridlines = False
    value_axis.visible = False
 
  elif n==15:  
    slide = tsg_ppt.slides[3]  
    chart_data = ChartData()
    chart_data.categories = ['1', '2', '3']
    chart_data.add_series('1', (a[0], b[0], c[0], d[0], e[0], f[0], ii[0], ji[0], k[0], l[0], m[0], nn[0], o[0], p[0], q[0]))
    chart_data.add_series('2', (a[1], b[1], c[1], d[1], e[1], f[1], ii[1], ji[1], k[1], l[1], m[1], nn[1], o[1], p[1], q[1]))
    chart_data.add_series('3', (a[2], b[2], c[2], d[2], e[2], f[2], ii[2], ji[2], k[2], l[2], m[2], nn[2], o[2], p[2], q[2]))
    x,y,cx,cy = Inches(0.3), Inches(2.25), Cm(14.14), Cm(11.39)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    value_axis = chart.value_axis
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.font.bold = True
    bar_plot = chart.plots[0]
    bar_plot.gap_width = 13
    chart.has_legend = False
    data_labels.font.size = Pt(12)
    data_labels.number_format = '0'
    data_labels.font.color.rgb = RGBColor(0, 0, 0)
    category_axis = chart.category_axis
    category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    category_axis.tick_labels.font.size = Pt(12)
    value_axis.minor_tick_mark = XL_TICK_MARK.NONE
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    value_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.has_major_gridlines = False
    value_axis.has_major_gridlines = False
    value_axis.visible = False
  elif n==16:  
    slide = tsg_ppt.slides[3]  
    chart_data = ChartData()
    chart_data.categories = ['1', '2', '3']
    chart_data.add_series('1', (a[0], b[0], c[0], d[0], e[0], f[0], ii[0], ji[0], k[0], l[0], m[0], nn[0], o[0], p[0], q[0], r[0]))
    chart_data.add_series('2', (a[1], b[1], c[1], d[1], e[1], f[1], ii[1], ji[1], k[1], l[1], m[1], nn[1], o[1], p[1], q[1], r[1]))
    chart_data.add_series('3', (a[2], b[2], c[2], d[2], e[2], f[2], ii[2], ji[2], k[2], l[2], m[2], nn[2], o[2], p[2], q[2], r[2]))
    x,y,cx,cy = Inches(0.3), Inches(2.25), Cm(14.14), Cm(11.39)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    value_axis = chart.value_axis
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.font.bold = True
    bar_plot = chart.plots[0]
    bar_plot.gap_width = 13
    chart.has_legend = False
    data_labels.font.size = Pt(12)
    data_labels.number_format = '0'
    data_labels.font.color.rgb = RGBColor(0, 0, 0)
    category_axis = chart.category_axis
    category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    category_axis.tick_labels.font.size = Pt(12)
    value_axis.minor_tick_mark = XL_TICK_MARK.NONE
    category_axis.major_tick_mark = XL_TICK_MARK.NONE
    value_axis.major_tick_mark = XL_TICK_MARK.NONE
    category_axis.has_major_gridlines = False
    value_axis.has_major_gridlines = False
    value_axis.visible = False
  for i in range(0, n):
    be=10
    if n==3:
      t=7.34
      he=3.55
      if diff[i] > 0:
        img_path='zoldnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Inches (0.425)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(0.39)
        table.cell(0, 0).text = "+" + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
	paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      elif diff[i] < 0:
        img_path='pirosnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
	pic = slide.shapes.add_picture(img_path, left, top, height=height)
      	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Inches (0.425)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(0.39)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      else:
        img_path='keknyil.png'
        left = Inches(6.25)
        height = Cm(0.43)
        top = Cm(t+0.33+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
       	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t-0.02+i*he)
        width = Inches(0.8)
        height = Inches (0.425)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(0.39)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
    elif n==5:
      t=6.66
      he=2.12
      if diff[i] > 0:
        img_path='zoldnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Inches (0.425)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(0.39)
        table.cell(0, 0).text = "+" + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      elif diff[i] < 0:
        img_path='pirosnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
	pic = slide.shapes.add_picture(img_path, left, top, height=height)
      	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Inches (0.425)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(0.39)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      else:
        img_path='keknyil.png'
        left = Inches(6.25)
        height = Cm(0.43)
        top = Cm(t+0.2+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
       	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Inches (0.425)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(0.39)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
    elif n==8:
      t=6.26
      he=1.32
      if diff[i] > 0:
        img_path='zoldnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Inches (0.425)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(0.39)
        table.cell(0, 0).text = "+" + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      elif diff[i] < 0:
        img_path='pirosnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
	pic = slide.shapes.add_picture(img_path, left, top, height=height)
      	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Inches (0.425)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(0.39)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      else:
        img_path='keknyil.png'
        left = Inches(6.25)
        height = Cm(0.43)
        top = Cm(t+0.2+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
       	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Inches (0.425)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(0.39)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
    elif n==9:
      t=6.18
      he=1.18
      if diff[i] > 0:
        img_path='zoldnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Inches (0.425)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(0.39)
        table.cell(0, 0).text = "+" + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      elif diff[i] < 0:
        img_path='pirosnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
	pic = slide.shapes.add_picture(img_path, left, top, height=height)
      	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Inches (0.425)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(0.39)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      else:
        img_path='keknyil.png'
        left = Inches(6.25)
        height = Cm(0.43)
        top = Cm(t+0.2+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
       	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Inches (0.425)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(0.39)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
    elif n==11:
      t=6.15
      he=0.97
      if diff[i] > 0:
        img_path='zoldnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Cm(0.9)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(1)
        table.cell(0, 0).text = "+" + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      elif diff[i] < 0:
        img_path='pirosnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
	pic = slide.shapes.add_picture(img_path, left, top, height=height)
      	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Cm(0.9)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(1)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      else:
        img_path='keknyil.png'
        left = Inches(6.25)
        height = Cm(0.43)
        top = Cm(t+0.2+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
       	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Cm(0.9)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(1)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
    elif n==13:
      t=6.15
      he=0.815
      if diff[i] > 0:
        img_path='zoldnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Cm(0.8)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(1)
        table.cell(0, 0).text = "+" + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      elif diff[i] < 0:
        img_path='pirosnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
	pic = slide.shapes.add_picture(img_path, left, top, height=height)
      	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Cm(0.8)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(1)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      else:
        img_path='keknyil.png'
        left = Inches(6.25)
        height = Cm(0.43)
        top = Cm(t+0.2+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
       	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Cm(0.8)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(1)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
    elif n==14:
      be=10
      ch=0.68
      t=6.15
      he=0.76
      if diff[i] > 0:
        img_path='zoldnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Cm(ch)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(1)
        table.cell(0, 0).text = "+" + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      elif diff[i] < 0:
        img_path='pirosnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
	pic = slide.shapes.add_picture(img_path, left, top, height=height)
      	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Cm(ch)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(1)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      else:
        img_path='keknyil.png'
        left = Inches(6.25)
        height = Cm(0.43)
        top = Cm(t+0.2+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
       	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Cm(ch)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(1)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
    elif n==15:
      be=10
      ch=0.68
      t=6.12
      he=0.71
      if diff[i] > 0:
        img_path='zoldnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Cm(ch)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(1)
        table.cell(0, 0).text = "+" + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      elif diff[i] < 0:
        img_path='pirosnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
	pic = slide.shapes.add_picture(img_path, left, top, height=height)
      	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Cm(ch)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(1)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      else:
        img_path='keknyil.png'
        left = Inches(6.25)
        height = Cm(0.43)
        top = Cm(t+0.2+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
       	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Cm(ch)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(1)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
    elif n==16:
      be=10
      ch=0.68
      t=6.1
      he=0.665
      if diff[i] > 0:
        img_path='zoldnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Cm(ch)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(1)
        table.cell(0, 0).text = "+" + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      elif diff[i] < 0:
        img_path='pirosnyil.png'
        left = Inches(6.25)
        height = Cm(0.6)
        top = Cm(t+i*he)
	pic = slide.shapes.add_picture(img_path, left, top, height=height)
      	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Cm(ch)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(1)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
      else:
        img_path='keknyil.png'
        left = Inches(6.25)
        height = Cm(0.43)
        top = Cm(t+0.2+i*he)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
       	rows=1
        cols=1
        left = Inches(5.83)
        top = Cm(t+i*he)
        width = Inches(0.8)
        height = Cm(ch)# set column widths
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Cm(1)
        table.cell(0, 0).text = " " + str(diff[i])
  	cell = table.rows[0].cells[0]
        paragraph = cell.textframe.paragraphs[0]
        paragraph.font.size = Pt(be)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.horizontal_anchor = MSO_ANCHOR.MIDDLE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(124,124,124)
  return tsg_ppt   

def fill_slide_4_to_10(tsg_ppt, a, n):
  slide = tsg_ppt.slides[a]
  chart_data = ChartData()
  chart_data.categories = ['1', '2', '3']
  chart_data.add_series('1', (a[0], b[0], c[0], d[0], e[0], f[0], ii[0], ji[0], k[0]))
  chart_data.add_series('2', (a[1], b[1], c[1], d[1], e[1], f[1], ii[1], ji[1], k[1]))
  chart_data.add_series('3', (a[2], b[2], c[2], d[2], e[2], f[2], ii[2], ji[2], k[2]))
  x,y,cx,cy = Inches(0.3), Inches(2.25), Cm(14.14), Cm(11.39)
  graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data)
  chart = graphic_frame.chart
  value_axis = chart.value_axis
  chart.plots[0].has_data_labels = True
  data_labels = chart.plots[0].data_labels
  data_labels.font.bold = True
  bar_plot = chart.plots[0]
  bar_plot.gap_width = 13
  chart.has_legend = False
  data_labels.font.size = Pt(12)
  data_labels.number_format = '0'
  data_labels.font.color.rgb = RGBColor(0, 0, 0)
  category_axis = chart.category_axis
  category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
  category_axis.tick_labels.font.size = Pt(12)
  value_axis.minor_tick_mark = XL_TICK_MARK.NONE
  category_axis.major_tick_mark = XL_TICK_MARK.NONE
  value_axis.major_tick_mark = XL_TICK_MARK.NONE
  category_axis.has_major_gridlines = False
  value_axis.has_major_gridlines = False
  value_axis.visible = False
  return tsg_ppt
############MAIN#####################  
tsg_ppt=Presentation('tsg_templ_uj.pptx')
orglongname = "Verkauf Deutschland"
orgshortname = "VLD"
q1 = ("20", "30", "50")

h = 1.6
a = [30, 40, 30]
b = [20, 30, 50]
c = [20, 59.4, 20.6]
d = [34, 16, 50]
e = [10, 20, 70]
f = [15, 75, 10]
g = [34, 26, 10, 20 ,10 ]
h = [15, 24, 11, 30, 20 ]
ii = [18, 22, 70]
ji = [20, 30, 50]
k = [30,40,30]
l = [20,30, 50]
m = [20,59.4,20.6]
nn = [34, 16, 50]
o = [10, 20, 70]
p = [15, 75, 10]
q = [34, 26, 40]
r = [15, 24, 61]
s = [18, 22, 70]
org_names = ["TSG", "GF OG", "VLD", "RVL N", "RVL O", "RVL W", "RVL M", "RVL SW", "RVL S"] #tomb1
org_percents = (("23", "22", "55"), ("10", "40", "50"), ("25", "35", "40"))
org_participants = ("40", "42", "33", "23", "55", "15", "67", "89", "56")
n = 9#len(org_names)
print(n)
diff = [1, 1, 0, -1, 9, 2, -3, 2, 0, 1, -1, 5, -8, -14, 5, 0] 
#picture_insert(tsg_ppt, 3, n, h)
fill_slide_0_and_titles(tsg_ppt, orglongname)
fill_slide_1(tsg_ppt, org_names, org_participants)  
fill_slide_2(tsg_ppt, diff)
fill_slide_3(tsg_ppt, n)
fill_slide_4_to_10(tsg_ppt, 4, n)
tsg_ppt.save("outteszt/"+"Abkurzung_der_OrgEinheit_TSG_Leadership_Survey"+".pptx")

